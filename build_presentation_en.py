# -*- coding: utf-8 -*-
"""Generate a professional English (LTR) PayNest sales presentation (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x0A,0x23,0x3A); BRAND=RGBColor(0x0C,0x8C,0xE8); BRAND_DARK=RGBColor(0x07,0x52,0x8C)
BRAND_DEEP=RGBColor(0x04,0x2A,0x4A); LIGHT=RGBColor(0xEA,0xF5,0xFD); WHITE=RGBColor(0xFF,0xFF,0xFF)
SLATE=RGBColor(0x47,0x55,0x69); SLATE_LT=RGBColor(0x94,0xA3,0xB8); GREEN=RGBColor(0x10,0xB9,0x81)
AMBER=RGBColor(0xF5,0x9E,0x0B); VIOLET=RGBColor(0x7C,0x3A,0xED); ROSE=RGBColor(0xE1,0x1D,0x48)
CARD_BG=RGBColor(0xF6,0xF9,0xFC)
FONT="Segoe UI"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def rect(s,x,y,w,h,c,shape=MSO_SHAPE.RECTANGLE,line=None):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=c
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space_after=6,ls=1.0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_left=tf.margin_right=Pt(4); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space_after); p.line_spacing=ls
        for (t,sz,b,c,f) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=f
    return tb
def R(t,sz,b=False,c=NAVY,f=FONT): return (t,sz,b,c,f)
def card(s,x,y,w,h,c=WHITE,line=RGBColor(0xE2,0xE8,0xF0)):
    sp=rect(s,x,y,w,h,c,MSO_SHAPE.ROUNDED_RECTANGLE,line); sp.adjustments[0]=0.06; return sp
def chip(s,x,y,w,h,c):
    sp=rect(s,x,y,w,h,c,MSO_SHAPE.ROUNDED_RECTANGLE); sp.adjustments[0]=0.5; return sp
def footer(s,n):
    text(s,Inches(0.45),Inches(6.95),Inches(4),Inches(0.4),
         [[R("PayNest",12,True,BRAND),R("  •  Payroll & HR Platform",11,False,SLATE_LT)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,Inches(11.6),Inches(6.95),Inches(1.3),Inches(0.4),
         [[R(f"{n} / 15",11,False,SLATE_LT)]],align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
def header(s,title,kicker=None):
    rect(s,0,0,SW,Inches(0.18),BRAND); rect(s,Inches(0.6),Inches(0.55),Inches(0.16),Inches(0.7),BRAND)
    if kicker: text(s,Inches(0.9),Inches(0.5),Inches(7.45),Inches(0.4),[[R(kicker,13,True,BRAND)]])
    text(s,Inches(0.9),Inches(0.85),Inches(11),Inches(0.8),[[R(title,30,True,NAVY)]])
    rect(s,Inches(0.9),Inches(1.72),Inches(8.8),Pt(2),LIGHT)

# 1 COVER
s=slide(); bg(s,BRAND_DEEP)
rect(s,Inches(9.5),Inches(-1.2),Inches(5.5),Inches(5.5),RGBColor(0x0A,0x3E,0x68),MSO_SHAPE.OVAL)
rect(s,Inches(10.8),Inches(3.6),Inches(4.5),Inches(4.5),RGBColor(0x0A,0x46,0x78),MSO_SHAPE.OVAL)
s.shapes.add_picture("brand/paynest-icon.png",Inches(0.9),Inches(1.95),Inches(1.05),Inches(1.05))
text(s,Inches(2.1),Inches(2.1),Inches(5),Inches(0.8),
     [[R("Pay",30,True,WHITE),R("Nest",30,True,BRAND)]],anchor=MSO_ANCHOR.MIDDLE)
text(s,Inches(0.9),Inches(3.35),Inches(11),Inches(1.6),
     [[R("Your team. Your payroll.",44,True,WHITE)],[R("One complete system.",44,True,BRAND)]])
text(s,Inches(0.95),Inches(5.25),Inches(10),Inches(0.7),
     [[R("A cloud HR & payroll platform built for businesses in Jordan and the Gulf — attendance, payroll, leave, and reviews in one place.",17,False,RGBColor(0xC7,0xDD,0xF0))]])
chip(s,Inches(0.95),Inches(6.2),Inches(3.0),Inches(0.55),BRAND)
text(s,Inches(0.95),Inches(6.2),Inches(3.0),Inches(0.55),[[R("Company Overview",14,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# 2 OVERVIEW
s=slide(); bg(s,WHITE); header(s,"What is PayNest?","Overview")
text(s,Inches(0.9),Inches(2.1),Inches(5.9),Inches(2.4),
     [[R("PayNest is an all-in-one payroll and HR system that brings everything a company needs into one secure, easy-to-use platform.",18,False,SLATE)],
      [R("",8,False,SLATE)],
      [R("From uploading attendance, to auto-calculating salaries, to managing leave and reviews — it's all automated and accurate.",18,False,SLATE)]],ls=1.15)
items=[("⚙️","Full Automation","Salaries & deductions calculated automatically"),
       ("🔒","Secure Isolation","Each company's data fully isolated"),
       ("🌐","Arabic / English","Fully bilingual interface"),
       ("📊","Live Reporting","Real-time dashboard & insights")]
y=Inches(2.1)
for e,t1,t2 in items:
    card(s,Inches(6.9),y,Inches(5.7),Inches(1.0))
    text(s,Inches(7.1),y+Inches(0.12),Inches(0.8),Inches(0.8),[[R(e,26,False,NAVY)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,Inches(7.95),y+Inches(0.14),Inches(4.45),Inches(0.8),[[R(t1,17,True,NAVY)],[R(t2,13,False,SLATE)]])
    y+=Inches(1.13)
footer(s,2)

# 3 PROBLEM
s=slide(); bg(s,WHITE); header(s,"The Challenges Companies Face","The Problem")
probs=[("⏱️","Manual payroll","Hours of work and human error every month"),
       ("📑","Scattered data","Spreadsheets and documents everywhere"),
       ("🧾","Attendance & leave","Hard to link hours to pay"),
       ("⚖️","Social security & deductions","Complex local-law compliance"),
       ("🔎","No oversight","No clear record of changes"),
       ("👥","Access control","Who sees what? Who approves what?")]
xs=[Inches(0.7),Inches(6.85)]; y0=Inches(2.05)
for i,(e,t1,t2) in enumerate(probs):
    col=i%2; row=i//2; x=xs[col]; y=y0+row*Inches(1.55)
    card(s,x,y,Inches(5.78),Inches(1.4),CARD_BG)
    chip(s,x+Inches(0.25),y+Inches(0.25),Inches(0.62),Inches(0.62),RGBColor(0xFD,0xE8,0xE8))
    text(s,x+Inches(0.25),y+Inches(0.25),Inches(0.62),Inches(0.62),[[R(e,20,False,NAVY)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(1.05),y+Inches(0.22),Inches(4.6),Inches(1.0),[[R(t1,17,True,NAVY)],[R(t2,13,False,SLATE)]])
footer(s,3)

# 4 SOLUTION
s=slide(); bg(s,BRAND_DEEP); rect(s,0,0,SW,Inches(0.18),BRAND)
text(s,Inches(0.9),Inches(0.8),Inches(7.45),Inches(0.5),[[R("The Solution",15,True,BRAND)]])
text(s,Inches(0.9),Inches(1.35),Inches(11),Inches(1.6),
     [[R("PayNest unifies every process into",34,True,WHITE)],[R("one secure, fully automated platform",34,True,BRAND)]])
flow=[("📥","Upload data","Attendance via Excel"),("⚙️","Auto-calculate","Accurate pay & deductions"),
      ("✅","Approve & track","Clear leave & review flow"),("📊","Analyze & decide","Live reports & dashboard")]
x=Inches(0.7); cw=Inches(2.95); gap=Inches(0.18)
for i,(e,t1,t2) in enumerate(flow):
    cx=x+i*(cw+gap); c=rect(s,cx,Inches(3.6),cw,Inches(2.6),RGBColor(0x0A,0x3E,0x68),MSO_SHAPE.ROUNDED_RECTANGLE); c.adjustments[0]=0.08
    chip(s,cx+Inches(1.12),Inches(3.95),Inches(0.7),Inches(0.7),BRAND)
    text(s,cx+Inches(1.12),Inches(3.95),Inches(0.7),Inches(0.7),[[R(e,24,False,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,cx+Inches(0.1),Inches(4.85),cw-Inches(0.2),Inches(0.5),[[R(t1,18,True,WHITE)]],align=PP_ALIGN.CENTER)
    text(s,cx+Inches(0.1),Inches(5.35),cw-Inches(0.2),Inches(0.7),[[R(t2,13,False,RGBColor(0xBF,0xD8,0xEE))]],align=PP_ALIGN.CENTER)
text(s,Inches(11.6),Inches(6.95),Inches(1.3),Inches(0.4),[[R("4 / 15",11,False,SLATE_LT)]],align=PP_ALIGN.RIGHT)

# 5 FEATURES
s=slide(); bg(s,WHITE); header(s,"Key Features","What PayNest Offers")
feats=[("💰","Smart Payroll",BRAND,"Auto calc — hourly or daily mode"),
       ("🕒","Time & Attendance",GREEN,"Upload fingerprint, link to pay"),
       ("📅","Leave & Permissions",AMBER,"Request & approve with clear flow"),
       ("🧑‍🤝‍🧑","Org Structure",VIOLET,"Assign supervisors by drag & drop"),
       ("⭐","Employee Reviews",BRAND,"Periodic reviews by supervisors"),
       ("📁","Document Vault",ROSE,"Store each employee's certificates"),
       ("🛡️","Audit Log",NAVY,"Track every change in the system"),
       ("🏢","Multi-Company",GREEN,"Manage several companies, isolated")]
x0=Inches(0.6); y0=Inches(2.0); cw=Inches(3.0); ch=Inches(2.25); gx=Inches(0.13); gy=Inches(0.2)
for i,(e,t1,col,t2) in enumerate(feats):
    c=i%4; r=i//4; x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    card(s,x,y,cw,ch)
    chip(s,x+Inches(0.25),y+Inches(0.25),Inches(0.7),Inches(0.7),col)
    text(s,x+Inches(0.25),y+Inches(0.25),Inches(0.7),Inches(0.7),[[R(e,22,False,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(0.2),y+Inches(1.05),cw-Inches(0.4),Inches(0.5),[[R(t1,17,True,NAVY)]])
    text(s,x+Inches(0.2),y+Inches(1.5),cw-Inches(0.4),Inches(0.7),[[R(t2,12.5,False,SLATE)]],ls=1.1)
footer(s,5)

# 6 PAYROLL
s=slide(); bg(s,WHITE); header(s,"Smart Payroll","Flagship Feature")
text(s,Inches(0.9),Inches(2.0),Inches(6.0),Inches(0.6),[[R("Two flexible modes for your business:",18,True,NAVY)]])
card(s,Inches(0.9),Inches(2.65),Inches(6.0),Inches(1.35),LIGHT)
text(s,Inches(1.15),Inches(2.8),Inches(5.5),Inches(1.1),
     [[R("⏱️ Hourly Mode",17,True,BRAND_DARK)],[R("Computes hourly rate, deducts shortfalls and adds overtime automatically.",13.5,False,SLATE)]],ls=1.1)
card(s,Inches(0.9),Inches(4.15),Inches(6.0),Inches(1.35),LIGHT)
text(s,Inches(1.15),Inches(4.3),Inches(5.5),Inches(1.1),
     [[R("📆 Daily Mode",17,True,BRAND_DARK)],[R("Computes pay day-by-day with absence and paid leave.",13.5,False,SLATE)]],ls=1.1)
incs=[("Base salary & allowances",GREEN),("Shortfall & absence deductions",ROSE),
      ("Overtime hours",BRAND),("Social security automatically",VIOLET),
      ("Manual bonuses & deductions",AMBER),("Net pay in one click",NAVY)]
y=Inches(2.05)
for t1,col in incs:
    chip(s,Inches(7.3),y+Inches(0.06),Inches(0.22),Inches(0.22),col)
    text(s,Inches(7.7),y,Inches(5.3),Inches(0.45),[[R(t1,15.5,False,NAVY)]],anchor=MSO_ANCHOR.MIDDLE)
    y+=Inches(0.62)
footer(s,6)

# 7 ATTENDANCE
s=slide(); bg(s,WHITE); header(s,"Time & Attendance","Accurate Tracking")
text(s,Inches(0.9),Inches(2.05),Inches(5.9),Inches(2.6),
     [[R("Upload the monthly fingerprint file as Excel, and PayNest reads each employee's clock-in/out and links it straight to payroll.",18,False,SLATE)],
      [R("",8,False,SLATE)],
      [R("It auto-detects worked hours, shortfalls, overtime and absent days — per month and year.",18,False,SLATE)]],ls=1.18)
stats=[("Worked Hours","Total productive hours",GREEN),("Required Hours","Monthly target",BRAND),
       ("Hour Gap","Shortfall / overtime",AMBER),("Monthly History","Full attendance record",VIOLET)]
y=Inches(2.1)
for t1,t2,col in stats:
    card(s,Inches(6.9),y,Inches(5.7),Inches(1.0))
    rect(s,Inches(6.9),y,Inches(0.22),Inches(1.0),col,MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,Inches(7.25),y+Inches(0.16),Inches(5.0),Inches(0.7),[[R(t1,17,True,NAVY)],[R(t2,12.5,False,SLATE_LT)]])
    y+=Inches(1.13)
footer(s,7)

# 8 LEAVE
s=slide(); bg(s,WHITE); header(s,"Leave & Permissions","Clear Approval Flow")
steps=[("1","Employee submits","Leave or short permission (1/2/3 hrs)",BRAND),
       ("2","Supervisor reviews","Approve or reject — decision is final",AMBER),
       ("3","Instant notification","Employee is notified immediately",GREEN),
       ("4","HR oversight","View & archive, no extra step",VIOLET)]
y=Inches(2.0)
for num,t1,t2,col in steps:
    chip(s,Inches(0.9),y+Inches(0.18),Inches(0.7),Inches(0.7),col)
    text(s,Inches(0.9),y+Inches(0.18),Inches(0.7),Inches(0.7),[[R(num,22,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    card(s,Inches(1.85),y,Inches(10.6),Inches(1.05))
    text(s,Inches(2.15),y+Inches(0.16),Inches(10.0),Inches(0.8),[[R(t1,18,True,NAVY)],[R(t2,14,False,SLATE)]])
    y+=Inches(1.2)
text(s,Inches(0.9),Inches(6.35),Inches(11.5),Inches(0.5),
     [[R("Bonus: attach a file with the request (medical note / chat screenshot); supports annual, sick and unpaid leave.",13.5,True,BRAND_DARK)]],align=PP_ALIGN.CENTER)
footer(s,8)

# 9 ORG + EVAL
s=slide(); bg(s,WHITE); header(s,"Org Structure & Reviews","Team Management")
card(s,Inches(0.7),Inches(2.0),Inches(6.05),Inches(4.4),CARD_BG)
text(s,Inches(0.95),Inches(2.2),Inches(5.55),Inches(0.6),[[R("🧑‍🤝‍🧑  Org Structure",19,True,NAVY)]])
for i,t in enumerate(["Assign supervisors via interactive drag-and-drop board",
                      "Support multiple supervisors per employee",
                      "Automatic cycle prevention",
                      "Each employee sees their manager and reports"]):
    text(s,Inches(0.95),Inches(2.85)+i*Inches(0.62),Inches(5.55),Inches(0.55),[[R("•  ",15,True,BRAND),R(t,14,False,SLATE)]],ls=1.05)
card(s,Inches(6.95),Inches(2.0),Inches(5.7),Inches(4.4),CARD_BG)
text(s,Inches(7.2),Inches(2.2),Inches(5.2),Inches(0.6),[[R("⭐  Employee Reviews",19,True,NAVY)]])
for i,t in enumerate(["Periodic reviews by supervisors and HR",
                      "Reviews tied to the org structure",
                      "Complete performance record per employee",
                      "A fair basis for bonuses and promotions"]):
    text(s,Inches(7.2),Inches(2.85)+i*Inches(0.62),Inches(5.2),Inches(0.55),[[R("•  ",15,True,AMBER),R(t,14,False,SLATE)]],ls=1.05)
footer(s,9)

# 10 ROLES
s=slide(); bg(s,WHITE); header(s,"Roles & Permissions","Everyone Sees Their Own")
roles=[("👑","Super Admin","Platform","Manages companies, plans & subscriptions",NAVY),
       ("🏢","Owner","Company","Full control of their company & settings",BRAND),
       ("👔","HR","Human Resources","Manages employees, payroll & reports",VIOLET),
       ("🧑‍💼","Employee","Self-service","Own portal: pay, leave, tasks",GREEN)]
x0=Inches(0.6); cw=Inches(3.0); gx=Inches(0.13)
for i,(e,t1,t2,t3,col) in enumerate(roles):
    x=x0+i*(cw+gx); card(s,x,Inches(2.15),cw,Inches(4.1))
    rect(s,x,Inches(2.15),cw,Inches(0.16),col,MSO_SHAPE.ROUNDED_RECTANGLE)
    chip(s,x+cw/2-Inches(0.55),Inches(2.65),Inches(1.1),Inches(1.1),col)
    text(s,x+cw/2-Inches(0.55),Inches(2.65),Inches(1.1),Inches(1.1),[[R(e,34,False,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(0.15),Inches(4.0),cw-Inches(0.3),Inches(0.6),[[R(t1,17,True,NAVY)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.15),Inches(4.5),cw-Inches(0.3),Inches(0.4),[[R(t2,13,True,col)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.25),Inches(5.0),cw-Inches(0.5),Inches(1.1),[[R(t3,13.5,False,SLATE)]],align=PP_ALIGN.CENTER,ls=1.12)
footer(s,10)

# 11 SECURITY
s=slide(); bg(s,BRAND_DEEP); rect(s,0,0,SW,Inches(0.18),BRAND)
text(s,Inches(0.9),Inches(0.7),Inches(7.45),Inches(0.5),[[R("Security",15,True,BRAND)]])
text(s,Inches(0.9),Inches(1.2),Inches(11),Inches(0.8),[[R("Enterprise-grade security 🛡️",30,True,WHITE)]])
secs=[("🔐","Full data isolation","Each company's data fully separated — no cross leaks"),
      ("🔑","Password hashing","One-way bcrypt — passwords can't be recovered"),
      ("✅","Identity validation","Real email and valid phone per country"),
      ("📜","Complete audit log","Every add, edit and delete is logged"),
      ("👁️","Fine-grained access","Each role sees only what it should"),
      ("☁️","Reliable cloud","Backups and high availability")]
xs=[Inches(0.7),Inches(6.85)]; y0=Inches(2.2)
for i,(e,t1,t2) in enumerate(secs):
    col=i%2; row=i//2; x=xs[col]; y=y0+row*Inches(1.45)
    c=rect(s,x,y,Inches(5.78),Inches(1.3),RGBColor(0x0A,0x3E,0x68),MSO_SHAPE.ROUNDED_RECTANGLE); c.adjustments[0]=0.08
    text(s,x+Inches(0.3),y+Inches(0.3),Inches(0.65),Inches(0.65),[[R(e,22,False,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(1.05),y+Inches(0.2),Inches(4.6),Inches(0.95),[[R(t1,16,True,WHITE)],[R(t2,12,False,RGBColor(0xBF,0xD8,0xEE))]],ls=1.05)
text(s,Inches(11.6),Inches(6.95),Inches(1.3),Inches(0.4),[[R("11 / 15",11,False,SLATE_LT)]],align=PP_ALIGN.RIGHT)

# 12 WHY
s=slide(); bg(s,WHITE); header(s,"Why PayNest?","What Sets Us Apart")
why=[("🌍","Built for our region","Supports social security & local labor law in Jordan and the Gulf"),
     ("🗣️","Bilingual","Full Arabic and English with RTL support"),
     ("⚡","Simple & fast","Up and running in minutes, no complexity"),
     ("🧩","Flexible & scalable","Hourly or daily, single or multi-company"),
     ("💸","Affordable","Plans that grow with your team size"),
     ("🤝","Local support","A team that understands you and speaks your language")]
xs=[Inches(0.7),Inches(6.85)]; y0=Inches(2.05)
for i,(e,t1,t2) in enumerate(why):
    col=i%2; row=i//2; x=xs[col]; y=y0+row*Inches(1.5)
    card(s,x,y,Inches(5.78),Inches(1.35))
    chip(s,x+Inches(0.25),y+Inches(0.3),Inches(0.7),Inches(0.7),LIGHT)
    text(s,x+Inches(0.25),y+Inches(0.3),Inches(0.7),Inches(0.7),[[R(e,22,False,NAVY)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(1.05),y+Inches(0.22),Inches(4.6),Inches(1.0),[[R(t1,17,True,NAVY)],[R(t2,13,False,SLATE)]],ls=1.08)
footer(s,12)

# 13 PRICING
s=slide(); bg(s,WHITE); header(s,"Subscription Plans","Plans That Grow With You")
plans=[("Starter","Up to 50 employees",["All payroll & attendance","Leave & reviews","Email support"],BRAND,False),
       ("Growth","51 to 100 employees",["Everything in Starter","Extended reports & audit log","Priority support"],NAVY,True),
       ("Enterprise","100+ employees",["Multi-company","Flexible custom limits","Dedicated support"],VIOLET,False)]
x0=Inches(0.95); cw=Inches(3.75); gx=Inches(0.35)
for i,(name,cap,feats,col,hot) in enumerate(plans):
    x=x0+i*(cw+gx); h=Inches(4.5) if hot else Inches(4.2); yy=Inches(2.05) if hot else Inches(2.2)
    card(s,x,yy,cw,h,WHITE if not hot else LIGHT,line=col if hot else RGBColor(0xE2,0xE8,0xF0))
    rect(s,x,yy,cw,Inches(0.16),col,MSO_SHAPE.ROUNDED_RECTANGLE)
    if hot:
        chip(s,x+cw/2-Inches(0.85),yy,Inches(1.7),Inches(0.5),col)
        text(s,x+cw/2-Inches(0.85),yy,Inches(1.7),Inches(0.5),[[R("Most Popular",12,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(0.2),yy+Inches(0.55),cw-Inches(0.4),Inches(0.6),[[R(name,20,True,NAVY)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.2),yy+Inches(1.15),cw-Inches(0.4),Inches(0.5),[[R(cap,15,True,col)]],align=PP_ALIGN.CENTER)
    for j,f in enumerate(feats):
        text(s,x+Inches(0.35),yy+Inches(1.85)+j*Inches(0.62),cw-Inches(0.6),Inches(0.55),[[R("✓  ",14,True,GREEN),R(f,13.5,False,SLATE)]],ls=1.05)
text(s,Inches(0.7),Inches(6.75),Inches(11.9),Inches(0.4),
     [[R("Pricing scales with size and features — contact us for a tailored quote.",13,True,SLATE)]],align=PP_ALIGN.CENTER)
footer(s,13)

# 14 DEMO
s=slide(); bg(s,BRAND_DEEP); rect(s,0,0,SW,Inches(0.18),BRAND)
text(s,Inches(0.9),Inches(0.7),Inches(7.45),Inches(0.5),[[R("Live Demo",15,True,BRAND)]])
text(s,Inches(0.9),Inches(1.2),Inches(11),Inches(0.8),[[R("See PayNest in action 🎬",30,True,WHITE)]])
s.shapes.add_movie("brand/PayNest-Demo.mp4",Inches(3.69),Inches(2.2),Inches(5.95),Inches(3.35),
                   poster_frame_image="brand/video-poster.png",mime_type="video/mp4")
text(s,Inches(3.4),Inches(5.75),Inches(6.55),Inches(0.5),[[R("A quick tour inside the system — click to play",13,False,RGBColor(0xBF,0xD8,0xEE))]],align=PP_ALIGN.CENTER)
text(s,Inches(11.6),Inches(6.95),Inches(1.3),Inches(0.4),[[R("14 / 15",11,False,SLATE_LT)]],align=PP_ALIGN.RIGHT)

# 15 CTA
s=slide(); bg(s,BRAND_DEEP)
rect(s,Inches(-1.5),Inches(4.5),Inches(6),Inches(6),RGBColor(0x0A,0x3E,0x68),MSO_SHAPE.OVAL)
s.shapes.add_picture("brand/paynest-icon.png",Inches(5.79),Inches(1.3),Inches(1.75),Inches(1.75))
text(s,Inches(1.5),Inches(3.2),Inches(10.3),Inches(0.9),[[R("Ready to digitize HR at your company",32,True,WHITE)]],align=PP_ALIGN.CENTER)
text(s,Inches(1.5),Inches(4.2),Inches(10.3),Inches(0.6),[[R("Let's book a free demo and get started today",18,False,RGBColor(0xC7,0xDD,0xF0))]],align=PP_ALIGN.CENTER)
chip(s,Inches(4.85),Inches(5.2),Inches(3.6),Inches(0.7),BRAND)
text(s,Inches(4.85),Inches(5.2),Inches(3.6),Inches(0.7),[[R("Start your trial now",17,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
text(s,Inches(1.5),Inches(6.25),Inches(10.3),Inches(0.5),[[R("paynest-tau.vercel.app",15,True,BRAND)]],align=PP_ALIGN.CENTER)

prs.save("PayNest-Presentation-EN.pptx")
print("SAVED PayNest-Presentation-EN.pptx | slides:",len(prs.slides._sldIdLst))
