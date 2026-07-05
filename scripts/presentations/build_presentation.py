# -*- coding: utf-8 -*-
"""Generate a professional Arabic (RTL) PayNest sales presentation (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY       = RGBColor(0x0A, 0x23, 0x3A)
BRAND      = RGBColor(0x0C, 0x8C, 0xE8)
BRAND_DARK = RGBColor(0x07, 0x52, 0x8C)
BRAND_DEEP = RGBColor(0x04, 0x2A, 0x4A)
LIGHT      = RGBColor(0xEA, 0xF5, 0xFD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SLATE      = RGBColor(0x47, 0x55, 0x69)
SLATE_LT   = RGBColor(0x94, 0xA3, 0xB8)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)
ROSE       = RGBColor(0xE1, 0x1D, 0x48)
CARD_BG    = RGBColor(0xF6, 0xF9, 0xFC)

AR = "Segoe UI"   # renders Arabic well on Windows

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_rtl(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    pPr.set('algn', 'r')


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.RIGHT, rtl=True, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, bold, color, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if rtl:
            set_rtl(p)
        for (txt, size, bold, color, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb


def R(txt, size, bold=False, color=NAVY, font=AR):
    return (txt, size, bold, color, font)


def footer(s, n):
    text(s, Inches(0.45), Inches(6.95), Inches(3), Inches(0.4),
         [[R("PayNest", 12, True, BRAND), R("  •  منصّة الرواتب والموارد البشرية", 11, False, SLATE_LT)]],
         align=PP_ALIGN.LEFT, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(11.6), Inches(6.95), Inches(1.3), Inches(0.4),
         [[R(f"{n} / 15", 11, False, SLATE_LT)]], align=PP_ALIGN.RIGHT, rtl=False,
         anchor=MSO_ANCHOR.MIDDLE)


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(0.18), BRAND)
    rect(s, Inches(12.55), Inches(0.55), Inches(0.16), Inches(0.7), BRAND)
    if kicker:
        text(s, Inches(5.0), Inches(0.5), Inches(7.45), Inches(0.4),
             [[R(kicker, 13, True, BRAND)]])
    text(s, Inches(3.5), Inches(0.85), Inches(8.95), Inches(0.8),
         [[R(title, 30, True, NAVY)]])
    rect(s, Inches(3.6), Inches(1.72), Inches(8.8), Pt(2), LIGHT)


def card(s, x, y, w, h, color=WHITE, line=RGBColor(0xE2,0xE8,0xF0)):
    c = rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=line)
    c.adjustments[0] = 0.06
    return c


def chip(s, x, y, w, h, color):
    c = rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.5
    return c


# ============================================================ 1 COVER
s = slide(); bg(s, BRAND_DEEP)
rect(s, 0, 0, SW, SH, BRAND_DEEP)
# accent shapes
b1 = rect(s, Inches(9.5), Inches(-1.2), Inches(5.5), Inches(5.5), BRAND_DARK, MSO_SHAPE.OVAL)
b1.fill.fore_color.rgb = RGBColor(0x0A, 0x3E, 0x68)
b2 = rect(s, Inches(10.8), Inches(3.6), Inches(4.5), Inches(4.5), RGBColor(0x0A,0x46,0x78), MSO_SHAPE.OVAL)
# logo icon
s.shapes.add_picture("brand/paynest-icon.png", Inches(0.9), Inches(1.95), Inches(1.05), Inches(1.05))
text(s, Inches(2.1), Inches(2.1), Inches(5), Inches(0.8),
     [[R("Pay", 30, True, WHITE, AR), R("Nest", 30, True, BRAND, AR)]],
     align=PP_ALIGN.LEFT, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(2.5), Inches(3.35), Inches(10), Inches(1.4),
     [[R("الرواتب والموارد البشرية،", 46, True, WHITE)],
      [R("في نظام واحد متكامل.", 46, True, BRAND)]],
     align=PP_ALIGN.RIGHT)
text(s, Inches(3.5), Inches(5.25), Inches(9), Inches(0.7),
     [[R("حلٌّ سحابي مصمّم لشركات الأردن والخليج — حضور، رواتب، إجازات، وتقييم في مكان واحد.",
         17, False, RGBColor(0xC7,0xDD,0xF0))]])
chip(s, Inches(9.55), Inches(6.2), Inches(3.3), Inches(0.55), BRAND)
text(s, Inches(9.55), Inches(6.2), Inches(3.3), Inches(0.55),
     [[R("عرض تعريفي للشركات", 14, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 2 OVERVIEW
s = slide(); bg(s, WHITE)
header(s, "ما هو PayNest؟", "نظرة عامة")
text(s, Inches(6.7), Inches(2.1), Inches(5.9), Inches(2.2),
     [[R("PayNest هو نظام متكامل لإدارة الرواتب والموارد البشرية، يجمع كل ما تحتاجه الشركة في منصّة واحدة سهلة الاستخدام وآمنة.", 18, False, SLATE, AR)],
      [R("", 8, False, SLATE)],
      [R("من رفع بصمة الحضور، إلى احتساب الرواتب تلقائياً، إلى إدارة الإجازات والتقييمات — كل شيء مؤتمت وبدقّة.", 18, False, SLATE, AR)]],
     line_spacing=1.15)
items = [("⚙️", "أتمتة كاملة", "احتساب الرواتب والخصومات تلقائياً"),
         ("🔒", "أمان وعزل", "بيانات كل شركة معزولة تماماً"),
         ("🌐", "عربي / إنجليزي", "واجهة ثنائية اللغة بالكامل"),
         ("📊", "تقارير لحظية", "لوحة تحكم ومؤشرات فورية")]
y = Inches(2.1)
for emoji, t1, t2 in items:
    card(s, Inches(0.7), y, Inches(5.7), Inches(1.0))
    text(s, Inches(5.55), y+Inches(0.12), Inches(0.8), Inches(0.8),
         [[R(emoji, 26, False, NAVY)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.95), y+Inches(0.14), Inches(4.45), Inches(0.8),
         [[R(t1, 17, True, NAVY)], [R(t2, 13, False, SLATE)]])
    y += Inches(1.13)
footer(s, 2)

# ============================================================ 3 PROBLEM
s = slide(); bg(s, WHITE)
header(s, "التحديات التي تواجه الشركات", "المشكلة")
probs = [("⏱️", "احتساب يدوي للرواتب", "ساعات من العمل وأخطاء بشرية كل شهر"),
         ("📑", "تشتّت البيانات", "ملفات Excel ومستندات متفرّقة"),
         ("🧾", "متابعة الحضور والإجازات", "صعوبة ربط الساعات بالرواتب"),
         ("⚖️", "الضمان والخصومات", "تعقيد في تطبيق القوانين المحلية"),
         ("🔎", "غياب الرقابة", "لا يوجد سجل واضح للتغييرات"),
         ("👥", "إدارة الصلاحيات", "من يرى ماذا؟ ومن يوافق على ماذا؟")]
xs = [Inches(6.85), Inches(0.7)]
y0 = Inches(2.05)
for i, (e, t1, t2) in enumerate(probs):
    col = i % 2; row = i // 2
    x = xs[col]; y = y0 + row*Inches(1.55)
    card(s, x, y, Inches(5.78), Inches(1.4), CARD_BG)
    chip(s, x+Inches(5.0), y+Inches(0.25), Inches(0.62), Inches(0.62), RGBColor(0xFD,0xE8,0xE8))
    text(s, x+Inches(5.0), y+Inches(0.25), Inches(0.62), Inches(0.62),
         [[R(e, 20, False, NAVY)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.25), y+Inches(0.22), Inches(4.6), Inches(1.0),
         [[R(t1, 17, True, NAVY)], [R(t2, 13, False, SLATE)]])
footer(s, 3)

# ============================================================ 4 SOLUTION
s = slide(); bg(s, BRAND_DEEP)
rect(s, 0, 0, SW, Inches(0.18), BRAND)
text(s, Inches(5.0), Inches(0.8), Inches(7.45), Inches(0.5), [[R("الحل", 15, True, BRAND)]])
text(s, Inches(2.5), Inches(1.35), Inches(9.95), Inches(1.6),
     [[R("PayNest يجمع كل العمليات في", 34, True, WHITE)],
      [R("منصّة واحدة، آمنة، ومؤتمتة بالكامل", 34, True, BRAND)]])
flow = [("📥", "ارفع البيانات", "بصمة وحضور بصيغة Excel"),
        ("⚙️", "احسب تلقائياً", "رواتب وخصومات بدقّة"),
        ("✅", "وافق وتابع", "إجازات وتقييمات بمسار واضح"),
        ("📊", "حلّل وقرّر", "تقارير ولوحة تحكم لحظية")]
x = Inches(0.7); cw = Inches(2.95); gap = Inches(0.18)
for i, (e, t1, t2) in enumerate(flow):
    cx = x + i*(cw+gap)
    c = rect(s, cx, Inches(3.6), cw, Inches(2.6), RGBColor(0x0A,0x3E,0x68), MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.08
    chip(s, cx+Inches(1.12), Inches(3.95), Inches(0.7), Inches(0.7), BRAND)
    text(s, cx+Inches(1.12), Inches(3.95), Inches(0.7), Inches(0.7),
         [[R(e, 24, False, WHITE)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+Inches(0.1), Inches(4.85), cw-Inches(0.2), Inches(0.5),
         [[R(t1, 18, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, cx+Inches(0.1), Inches(5.35), cw-Inches(0.2), Inches(0.7),
         [[R(t2, 13, False, RGBColor(0xBF,0xD8,0xEE))]], align=PP_ALIGN.CENTER)
text(s, Inches(11.6), Inches(6.95), Inches(1.3), Inches(0.4),
     [[R("4 / 15", 11, False, SLATE_LT)]], align=PP_ALIGN.RIGHT, rtl=False)

# ============================================================ 5 FEATURES GRID
s = slide(); bg(s, WHITE)
header(s, "المزايا الرئيسية", "ماذا يقدّم PayNest")
feats = [("💰", "الرواتب الذكية", BRAND, "احتساب تلقائي بنظام الساعات أو اليومي"),
         ("🕒", "الحضور والانصراف", GREEN, "رفع البصمة وربطها بالرواتب"),
         ("📅", "الإجازات والأذونات", AMBER, "طلب وموافقة بمسار واضح"),
         ("🧑‍🤝‍🧑", "الهيكل الإداري", VIOLET, "تحديد المشرفين بالسحب والإفلات"),
         ("⭐", "تقييم الموظفين", BRAND, "تقييمات دورية من المشرفين"),
         ("📁", "أرشيف الوثائق", ROSE, "حفظ شهادات ومستندات كل موظف"),
         ("🛡️", "سجل التدقيق", NAVY, "تتبّع كل تغيير في النظام"),
         ("🏢", "تعدّد الشركات", GREEN, "إدارة عدة شركات بعزل تام")]
x0 = Inches(0.6); y0 = Inches(2.0); cw = Inches(3.0); ch = Inches(2.25); gx = Inches(0.13); gy = Inches(0.2)
for i, (e, t1, col, t2) in enumerate(feats):
    c = i % 4; r = i // 4
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    card(s, x, y, cw, ch)
    chip(s, x+cw-Inches(0.95), y+Inches(0.25), Inches(0.7), Inches(0.7), col)
    text(s, x+cw-Inches(0.95), y+Inches(0.25), Inches(0.7), Inches(0.7),
         [[R(e, 22, False, WHITE)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.2), y+Inches(1.05), cw-Inches(0.4), Inches(0.5),
         [[R(t1, 17, True, NAVY)]])
    text(s, x+Inches(0.2), y+Inches(1.5), cw-Inches(0.4), Inches(0.7),
         [[R(t2, 12.5, False, SLATE)]], line_spacing=1.1)
footer(s, 5)

# ============================================================ 6 PAYROLL
s = slide(); bg(s, WHITE)
header(s, "إدارة الرواتب الذكية", "الميزة الأبرز")
text(s, Inches(6.6), Inches(2.0), Inches(6.0), Inches(0.6),
     [[R("نظامان مرنان حسب طبيعة شركتك:", 18, True, NAVY)]])
card(s, Inches(6.6), Inches(2.65), Inches(6.0), Inches(1.35), LIGHT)
text(s, Inches(6.85), Inches(2.8), Inches(5.5), Inches(1.1),
     [[R("⏱️ نظام الساعات", 17, True, BRAND_DARK, AR)],
      [R("يحتسب أجر الساعة ويخصم نقص الساعات ويضيف الإضافي تلقائياً.", 13.5, False, SLATE, AR)]],
     line_spacing=1.1)
card(s, Inches(6.6), Inches(4.15), Inches(6.0), Inches(1.35), LIGHT)
text(s, Inches(6.85), Inches(4.3), Inches(5.5), Inches(1.1),
     [[R("📆 النظام اليومي", 17, True, BRAND_DARK, AR)],
      [R("يحتسب الراتب يوماً بيوم مع الغياب والإجازات المدفوعة.", 13.5, False, SLATE, AR)]],
     line_spacing=1.1)
incs = [("الراتب الأساسي والعلاوات", GREEN),
        ("خصم نقص الساعات والغياب", ROSE),
        ("الساعات الإضافية (Overtime)", BRAND),
        ("الضمان الاجتماعي تلقائياً", VIOLET),
        ("المكافآت والخصومات اليدوية", AMBER),
        ("صافي الراتب بضغطة واحدة", NAVY)]
y = Inches(2.05)
for t1, col in incs:
    chip(s, Inches(6.15), y+Inches(0.06), Inches(0.22), Inches(0.22), col)
    text(s, Inches(0.7), y, Inches(5.3), Inches(0.45),
         [[R(t1, 15.5, False, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)
footer(s, 6)

# ============================================================ 7 ATTENDANCE
s = slide(); bg(s, WHITE)
header(s, "الحضور والانصراف", "التتبّع الدقيق")
text(s, Inches(6.7), Inches(2.05), Inches(5.9), Inches(2.4),
     [[R("ارفع ملف البصمة الشهري بصيغة Excel، وسيقوم PayNest بقراءة ساعات الدخول والخروج لكل موظف وربطها مباشرة باحتساب الراتب.", 18, False, SLATE)],
      [R("", 8, False, SLATE)],
      [R("يحدّد النظام تلقائياً: الساعات المنجزة، النقص، الإضافي، وأيام الغياب — لكل شهر وسنة على حدة.", 18, False, SLATE)]],
     line_spacing=1.18)
stats = [("الساعات المنجزة", "Worked Hours", GREEN),
         ("الساعات المطلوبة", "Required Hours", BRAND),
         ("فرق الساعات", "Hour Gap", AMBER),
         ("سجل شهري كامل", "Monthly History", VIOLET)]
y = Inches(2.1)
for t1, t2, col in stats:
    card(s, Inches(0.7), y, Inches(5.7), Inches(1.0))
    rect(s, Inches(6.18), y, Inches(0.22), Inches(1.0), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.95), y+Inches(0.16), Inches(5.0), Inches(0.7),
         [[R(t1, 17, True, NAVY)], [R(t2, 12.5, False, SLATE_LT, AR)]])
    y += Inches(1.13)
footer(s, 7)

# ============================================================ 8 LEAVES WORKFLOW
s = slide(); bg(s, WHITE)
header(s, "الإجازات والأذونات", "مسار موافقة واضح")
steps = [("1", "الموظف يقدّم الطلب", "إجازة أو إذن مغادرة (ساعة/ساعتين/٣)", BRAND),
         ("2", "المشرف يراجع", "موافقة أو رفض — قراره نهائي", AMBER),
         ("3", "إشعار فوري", "يصل الموظف بالنتيجة مباشرة", GREEN),
         ("4", "الموارد البشرية تتابع", "اطّلاع وأرشفة دون تدخّل", VIOLET)]
y = Inches(2.0)
for num, t1, t2, col in steps:
    chip(s, Inches(11.7), y+Inches(0.18), Inches(0.7), Inches(0.7), col)
    text(s, Inches(11.7), y+Inches(0.18), Inches(0.7), Inches(0.7),
         [[R(num, 22, True, WHITE)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    card(s, Inches(0.7), y, Inches(10.8), Inches(1.05))
    text(s, Inches(1.0), y+Inches(0.16), Inches(10.2), Inches(0.8),
         [[R(t1, 18, True, NAVY)], [R(t2, 14, False, SLATE)]])
    y += Inches(1.2)
text(s, Inches(0.7), Inches(6.35), Inches(10.8), Inches(0.5),
     [[R("ميزة إضافية: إرفاق ملف مع الطلب (تقرير طبي / صورة محادثة) ودعم الإجازات السنوية والمرضية وبدون راتب.",
         13.5, True, BRAND_DARK)]], align=PP_ALIGN.CENTER)
footer(s, 8)

# ============================================================ 9 ORG + EVAL
s = slide(); bg(s, WHITE)
header(s, "الهيكل الإداري والتقييم", "إدارة الفريق")
card(s, Inches(6.55), Inches(2.0), Inches(6.05), Inches(4.4), CARD_BG)
text(s, Inches(6.8), Inches(2.2), Inches(5.55), Inches(0.6),
     [[R("🧑‍🤝‍🧑  الهيكل الإداري", 19, True, NAVY)]])
for i, t in enumerate(["تحديد المشرفين بالسحب والإفلات على لوحة تفاعلية",
                       "دعم أكثر من مشرف للموظف الواحد",
                       "منع التسلسل الدائري تلقائياً",
                       "كل موظف يرى مشرفه ومن يشرف عليهم"]):
    text(s, Inches(6.8), Inches(2.85)+i*Inches(0.62), Inches(5.55), Inches(0.55),
         [[R("• ", 15, True, BRAND), R(t, 14, False, SLATE)]], line_spacing=1.05)
card(s, Inches(0.7), Inches(2.0), Inches(5.7), Inches(4.4), CARD_BG)
text(s, Inches(0.95), Inches(2.2), Inches(5.2), Inches(0.6),
     [[R("⭐  تقييم الموظفين", 19, True, NAVY)]])
for i, t in enumerate(["تقييمات دورية من المشرفين والموارد البشرية",
                       "ربط التقييم بالهيكل الإداري",
                       "سجل أداء متكامل لكل موظف",
                       "أساس عادل للمكافآت والترقيات"]):
    text(s, Inches(0.95), Inches(2.85)+i*Inches(0.62), Inches(5.2), Inches(0.55),
         [[R("• ", 15, True, AMBER), R(t, 14, False, SLATE)]], line_spacing=1.05)
footer(s, 9)

# ============================================================ 10 ROLES
s = slide(); bg(s, WHITE)
header(s, "الأدوار والصلاحيات", "كل شخص يرى ما يخصّه")
roles = [("👑", "المدير العام للنظام", "Super Admin", "إدارة الشركات والباقات والاشتراكات", NAVY),
         ("🏢", "مالك الشركة", "Owner", "تحكّم كامل بشركته وإعداداتها", BRAND),
         ("👔", "الموارد البشرية", "HR", "إدارة الموظفين والرواتب والتقارير", VIOLET),
         ("🧑‍💼", "الموظف", "Employee", "بوابة خاصة: راتبه، إجازاته، مهامه", GREEN)]
x0 = Inches(0.6); cw = Inches(3.0); gx = Inches(0.13)
for i, (e, t1, t2, t3, col) in enumerate(roles):
    x = x0 + i*(cw+gx)
    card(s, x, Inches(2.15), cw, Inches(4.1))
    rect(s, x, Inches(2.15), cw, Inches(0.16), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    chip(s, x+cw/2-Inches(0.55), Inches(2.65), Inches(1.1), Inches(1.1), col)
    text(s, x+cw/2-Inches(0.55), Inches(2.65), Inches(1.1), Inches(1.1),
         [[R(e, 34, False, WHITE)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.15), Inches(4.0), cw-Inches(0.3), Inches(0.6),
         [[R(t1, 17, True, NAVY)]], align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.15), Inches(4.5), cw-Inches(0.3), Inches(0.4),
         [[R(t2, 13, True, col)]], align=PP_ALIGN.CENTER, rtl=False)
    text(s, x+Inches(0.25), Inches(5.0), cw-Inches(0.5), Inches(1.1),
         [[R(t3, 13.5, False, SLATE)]], align=PP_ALIGN.CENTER, line_spacing=1.12)
footer(s, 10)

# ============================================================ 11 SECURITY
s = slide(); bg(s, BRAND_DEEP)
rect(s, 0, 0, SW, Inches(0.18), BRAND)
text(s, Inches(5.0), Inches(0.7), Inches(7.45), Inches(0.5), [[R("الأمان", 15, True, BRAND)]])
text(s, Inches(4.0), Inches(1.2), Inches(8.45), Inches(0.8),
     [[R("أمان على مستوى المؤسسات 🛡️", 30, True, WHITE)]])
secs = [("🔐", "عزل تام للبيانات", "بيانات كل شركة معزولة بالكامل — لا تسرّب بين الشركات"),
        ("🔑", "تشفير كلمات السر", "تشفير أحادي الاتجاه (bcrypt) — لا يمكن استرجاعها"),
        ("✅", "تحقّق من الهوية", "بريد إلكتروني حقيقي ورقم هاتف صحيح لكل دولة"),
        ("📜", "سجل تدقيق كامل", "كل إضافة وتعديل وحذف مسجّل ومؤرّخ"),
        ("👁️", "صلاحيات دقيقة", "كل دور يصل لما يخصّه فقط"),
        ("☁️", "بنية سحابية موثوقة", "نسخ احتياطي وتوفّر عالٍ")]
xs = [Inches(6.85), Inches(0.7)]
y0 = Inches(2.2)
for i, (e, t1, t2) in enumerate(secs):
    col = i % 2; row = i // 2
    x = xs[col]; y = y0 + row*Inches(1.45)
    c = rect(s, x, y, Inches(5.78), Inches(1.3), RGBColor(0x0A,0x3E,0x68), MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.08
    text(s, x+Inches(5.0), y+Inches(0.3), Inches(0.65), Inches(0.65),
         [[R(e, 22, False, WHITE)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.25), y+Inches(0.2), Inches(4.6), Inches(0.95),
         [[R(t1, 16, True, WHITE)], [R(t2, 12, False, RGBColor(0xBF,0xD8,0xEE))]], line_spacing=1.05)
text(s, Inches(11.6), Inches(6.95), Inches(1.3), Inches(0.4),
     [[R("11 / 15", 11, False, SLATE_LT)]], align=PP_ALIGN.RIGHT, rtl=False)

# ============================================================ 12 WHY US
s = slide(); bg(s, WHITE)
header(s, "لماذا PayNest؟", "ما يميّزنا")
why = [("🌍", "مصمّم لمنطقتنا", "يدعم الضمان الاجتماعي وقوانين العمل في الأردن والخليج"),
       ("🗣️", "ثنائي اللغة", "عربي وإنجليزي بالكامل مع دعم RTL"),
       ("⚡", "سهل وسريع", "تشغيل خلال دقائق دون تعقيد"),
       ("🧩", "مرن وقابل للتوسّع", "نظام ساعات أو يومي، شركة أو عدة شركات"),
       ("💸", "بأسعار مناسبة", "باقات تنمو مع حجم فريقك"),
       ("🤝", "دعم محلي", "فريق يفهم احتياجك ويتحدّث لغتك")]
xs = [Inches(6.85), Inches(0.7)]
y0 = Inches(2.05)
for i, (e, t1, t2) in enumerate(why):
    col = i % 2; row = i // 2
    x = xs[col]; y = y0 + row*Inches(1.5)
    card(s, x, y, Inches(5.78), Inches(1.35))
    chip(s, x+Inches(5.0), y+Inches(0.3), Inches(0.7), Inches(0.7), LIGHT)
    text(s, x+Inches(5.0), y+Inches(0.3), Inches(0.7), Inches(0.7),
         [[R(e, 22, False, NAVY)]], align=PP_ALIGN.CENTER, rtl=False, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.25), y+Inches(0.22), Inches(4.6), Inches(1.0),
         [[R(t1, 17, True, NAVY)], [R(t2, 13, False, SLATE)]], line_spacing=1.08)
footer(s, 12)

# ============================================================ 13 PRICING
s = slide(); bg(s, WHITE)
header(s, "باقات الاشتراك", "خطط تنمو معك")
plans = [("الباقة الأساسية", "حتى 50 موظف", ["كل مزايا الرواتب والحضور",
          "إدارة الإجازات والتقييم", "دعم عبر البريد"], BRAND, False),
         ("الباقة المتقدّمة", "من 51 إلى 100 موظف", ["كل مزايا الأساسية",
          "تقارير موسّعة وسجل تدقيق", "دعم ذو أولوية"], NAVY, True),
         ("باقة المؤسسات", "أكثر من 100 موظف", ["تعدّد الشركات",
          "حدود مرنة حسب الطلب", "دعم مخصّص"], VIOLET, False)]
x0 = Inches(0.95); cw = Inches(3.75); gx = Inches(0.35)
for i, (name, cap, feats, col, hot) in enumerate(plans):
    x = x0 + i*(cw+gx)
    h = Inches(4.5) if hot else Inches(4.2)
    yy = Inches(2.05) if hot else Inches(2.2)
    card(s, x, yy, cw, h, WHITE if not hot else LIGHT,
         line=col if hot else RGBColor(0xE2,0xE8,0xF0))
    rect(s, x, yy, cw, Inches(0.16), col, MSO_SHAPE.ROUNDED_RECTANGLE)
    if hot:
        chip(s, x+cw/2-Inches(0.85), yy-Inches(0.0), Inches(1.7), Inches(0.5), col)
        text(s, x+cw/2-Inches(0.85), yy, Inches(1.7), Inches(0.5),
             [[R("الأكثر طلباً", 12, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.2), yy+Inches(0.55), cw-Inches(0.4), Inches(0.6),
         [[R(name, 20, True, NAVY)]], align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.2), yy+Inches(1.15), cw-Inches(0.4), Inches(0.5),
         [[R(cap, 15, True, col)]], align=PP_ALIGN.CENTER)
    for j, f in enumerate(feats):
        text(s, x+Inches(0.3), yy+Inches(1.85)+j*Inches(0.62), cw-Inches(0.6), Inches(0.55),
             [[R("✓  ", 14, True, GREEN), R(f, 13.5, False, SLATE)]], align=PP_ALIGN.RIGHT, line_spacing=1.05)
text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.4),
     [[R("الأسعار تُحدّد حسب الحجم والمزايا — تواصل معنا لعرض مخصّص لشركتك.", 13, True, SLATE)]],
     align=PP_ALIGN.CENTER)
footer(s, 13)

# ============================================================ 14 DEMO / VIDEO
s = slide(); bg(s, BRAND_DEEP)
rect(s, 0, 0, SW, Inches(0.18), BRAND)
text(s, Inches(5.0), Inches(0.7), Inches(7.45), Inches(0.5), [[R("عرض حيّ", 15, True, BRAND)]])
text(s, Inches(3.5), Inches(1.2), Inches(8.95), Inches(0.8),
     [[R("شاهد PayNest أثناء العمل 🎬", 30, True, WHITE)]])
# embedded product-tour video
s.shapes.add_movie("brand/PayNest-Demo.mp4", Inches(3.69), Inches(2.2), Inches(5.95), Inches(3.35),
                   poster_frame_image="brand/video-poster.png", mime_type="video/mp4")
text(s, Inches(3.4), Inches(5.75), Inches(6.55), Inches(0.5),
     [[R("جولة سريعة داخل النظام — اضغط للتشغيل", 13, False, RGBColor(0xBF,0xD8,0xEE))]],
     align=PP_ALIGN.CENTER)
text(s, Inches(11.6), Inches(6.95), Inches(1.3), Inches(0.4),
     [[R("14 / 15", 11, False, SLATE_LT)]], align=PP_ALIGN.RIGHT, rtl=False)

# ============================================================ 15 CONTACT / CTA
s = slide(); bg(s, BRAND_DEEP)
rect(s, 0, 0, SW, SH, BRAND_DEEP)
b = rect(s, Inches(-1.5), Inches(4.5), Inches(6), Inches(6), RGBColor(0x0A,0x3E,0x68), MSO_SHAPE.OVAL)
s.shapes.add_picture("brand/paynest-icon.png", Inches(5.77), Inches(1.3), Inches(1.75), Inches(1.75))
text(s, Inches(2.5), Inches(3.2), Inches(8.3), Inches(0.9),
     [[R("جاهزون لرقمنة الموارد البشرية في شركتك", 32, True, WHITE)]], align=PP_ALIGN.CENTER)
text(s, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.6),
     [[R("لنحجز لك عرضاً توضيحياً مجانياً ونبدأ اليوم", 18, False, RGBColor(0xC7,0xDD,0xF0))]],
     align=PP_ALIGN.CENTER)
chip(s, Inches(4.85), Inches(5.2), Inches(3.6), Inches(0.7), BRAND)
text(s, Inches(4.85), Inches(5.2), Inches(3.6), Inches(0.7),
     [[R("ابدأ تجربتك الآن", 17, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(2.5), Inches(6.25), Inches(8.3), Inches(0.5),
     [[R("paynest-tau.vercel.app", 15, True, BRAND)]], align=PP_ALIGN.CENTER, rtl=False)

prs.save("PayNest-Presentation.pptx")
print("SAVED PayNest-Presentation.pptx  | slides:", len(prs.slides._sldIdLst))
