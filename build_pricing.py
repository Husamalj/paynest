# -*- coding: utf-8 -*-
"""Generate an attractive Arabic (RTL) PayNest pricing presentation (.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY=RGBColor(0x0A,0x23,0x3A); BRAND=RGBColor(0x0C,0x8C,0xE8); BRAND_DARK=RGBColor(0x07,0x52,0x8C)
BRAND_DEEP=RGBColor(0x04,0x2A,0x4A); LIGHT=RGBColor(0xEA,0xF5,0xFD); WHITE=RGBColor(0xFF,0xFF,0xFF)
SLATE=RGBColor(0x47,0x55,0x69); SLATE_LT=RGBColor(0x94,0xA3,0xB8); GREEN=RGBColor(0x10,0xB9,0x81)
AMBER=RGBColor(0xF5,0x9E,0x0B); GOLD=RGBColor(0xD9,0xA4,0x06); VIOLET=RGBColor(0x7C,0x3A,0xED)
F="Segoe UI"
ICON=os.path.join("brand","paynest-icon.png")

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c
def rect(s,x,y,w,h,c,shape=MSO_SHAPE.RECTANGLE,line=None,lw=1.0):
    sp=s.shapes.add_shape(shape,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=c
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def card(s,x,y,w,h,c=WHITE,line=RGBColor(0xE2,0xE8,0xF0),lw=1.0):
    sp=rect(s,x,y,w,h,c,MSO_SHAPE.ROUNDED_RECTANGLE,line,lw); sp.adjustments[0]=0.05; return sp
def chip(s,x,y,w,h,c):
    sp=rect(s,x,y,w,h,c,MSO_SHAPE.ROUNDED_RECTANGLE); sp.adjustments[0]=0.5; return sp
def text(s,x,y,w,h,runs,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.TOP,sa=4,ls=1.0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor; tf.margin_left=tf.margin_right=Pt(4); tf.margin_top=tf.margin_bottom=Pt(2)
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sa); p.line_spacing=ls
        pPr=p._p.get_or_add_pPr(); pPr.set("rtl","1")
        for (t,sz,b,c) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=F
    return tb
def R(t,sz,b=False,c=NAVY): return (t,sz,b,c)

# ───────────────────────── SLIDE 1: PRICING ─────────────────────────
s=slide(); bg(s,RGBColor(0xF6,0xF9,0xFC))
rect(s,0,0,SW,Inches(0.18),BRAND)
if os.path.exists(ICON): s.shapes.add_picture(ICON,Inches(11.9),Inches(0.45),Inches(0.9),Inches(0.9))
text(s,Inches(2.5),Inches(0.55),Inches(9.2),Inches(0.7),[[R("باقات وأسعار PayNest",34,True,NAVY)]])
text(s,Inches(2.5),Inches(1.25),Inches(9.2),Inches(0.5),[[R("نظام موارد بشرية ورواتب متكامل — بسعر يناسب حجم شركتك",16,False,SLATE)]])

tiers=[
    {"name":"الأساسية","sub":"حتى 50 موظف","per":"2 دينار / موظف / شهر","price":"≈ 350","unit":"دينار سنوياً","color":BRAND,"hot":False,
     "feats":["الرواتب (يومي + بالساعة)","الحضور والإجازات والمغادرات","التقييم والمهام","دعم عبر البريد"]},
    {"name":"المتقدّمة","sub":"من 51 إلى 100 موظف","per":"1.5 دينار / موظف / شهر","price":"≈ 890","unit":"دينار سنوياً","color":NAVY,"hot":True,
     "feats":["كل مزايا الأساسية","السلف والمكافآت والخصومات","الهيكل الإداري + الخطابات الرسمية","سجل تدقيق + تقارير","دعم ذو أولوية"]},
    {"name":"المؤسسات","sub":"أكثر من 100 موظف","per":"تسعير مخصّص","price":"تواصل","unit":"معنا","color":VIOLET,"hot":False,
     "feats":["كل مزايا المتقدّمة","تعدّد الشركات","حدود مرنة + تكامل خاص","دعم مخصّص"]},
]
# order right-to-left for RTL feel: Enterprise, Growth(center), Starter
x0=Inches(0.85); cw=Inches(3.75); gx=Inches(0.35)
positions=[x0+i*(cw+gx) for i in range(3)]
for i,t in enumerate(tiers):
    x=positions[i]
    h=Inches(5.15) if t["hot"] else Inches(4.8)
    yy=Inches(2.0) if t["hot"] else Inches(2.2)
    c=card(s,x,yy,cw,h,WHITE if not t["hot"] else RGBColor(0x07,0x2A,0x44),line=t["color"] if t["hot"] else RGBColor(0xE2,0xE8,0xF0),lw=2 if t["hot"] else 1)
    rect(s,x,yy,cw,Inches(0.16),t["color"],MSO_SHAPE.ROUNDED_RECTANGLE)
    txtcol = WHITE if t["hot"] else NAVY
    subcol = RGBColor(0xC7,0xDD,0xF0) if t["hot"] else SLATE_LT
    if t["hot"]:
        ch=chip(s,x+cw/2-Inches(0.95),yy-Inches(0.02),Inches(1.9),Inches(0.5),AMBER)
        text(s,x+cw/2-Inches(0.95),yy-Inches(0.02),Inches(1.9),Inches(0.5),[[R("⭐ الأكثر طلباً",12,True,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(0.2),yy+Inches(0.5),cw-Inches(0.4),Inches(0.5),[[R(t["name"],22,True,txtcol)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.2),yy+Inches(1.05),cw-Inches(0.4),Inches(0.4),[[R(t["sub"],13,False,subcol)]],align=PP_ALIGN.CENTER)
    # price
    text(s,x+Inches(0.2),yy+Inches(1.5),cw-Inches(0.4),Inches(0.7),[[R(t["price"],40,True,t["color"] if not t["hot"] else WHITE)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.2),yy+Inches(2.25),cw-Inches(0.4),Inches(0.35),[[R(t["unit"],13,True,subcol)]],align=PP_ALIGN.CENTER)
    text(s,x+Inches(0.2),yy+Inches(2.6),cw-Inches(0.4),Inches(0.3),[[R(t["per"],11,False,subcol)]],align=PP_ALIGN.CENTER)
    # features
    for j,fdef in enumerate(t["feats"]):
        text(s,x+Inches(0.3),yy+Inches(3.05)+j*Inches(0.36),cw-Inches(0.6),Inches(0.35),
             [[R(fdef+"  ",12,False,(RGBColor(0xE6,0xEF,0xF7) if t["hot"] else SLATE)),R("✓",12,True,GREEN)]],align=PP_ALIGN.RIGHT,ls=1.0)
text(s,Inches(0.85),Inches(7.0),Inches(11.6),Inches(0.4),
     [[R("الأسعار شاملة كل المزايا • خصم 15% عند الدفع السنوي • رسوم تجهيز لمرة واحدة 50–150 دينار",12,True,SLATE)]],align=PP_ALIGN.CENTER)

# ───────────────────────── SLIDE 2: VALUE / WHY ─────────────────────────
s=slide(); bg(s,BRAND_DEEP)
rect(s,0,0,SW,Inches(0.18),BRAND)
if os.path.exists(ICON): s.shapes.add_picture(ICON,Inches(0.85),Inches(0.55),Inches(0.95),Inches(0.95))
text(s,Inches(2.0),Inches(0.7),Inches(10.5),Inches(0.8),[[R("ليش PayNest يستحق؟",32,True,WHITE)]])
text(s,Inches(2.0),Inches(1.45),Inches(10.5),Inches(0.5),[[R("أوفر، أبسط، ومصمّم لشركات المنطقة",17,False,RGBColor(0xC7,0xDD,0xF0))]])

points=[
    ("💸","أرخص من موظف إداري","نظام كامل بأقل من تكلفة نصف راتب موظف بالسنة"),
    ("⚡","تشغيل خلال دقائق","بدون أجهزة معقّدة — رفع Excel وخلصت"),
    ("🌍","مصمّم لنا","ضمان اجتماعي، عربي/إنجليزي، قوانين المنطقة"),
    ("🔒","آمن ومعزول","بيانات كل شركة معزولة + سجل تدقيق كامل"),
    ("🧩","ينمو معك","من شركة صغيرة لعدة شركات بنفس النظام"),
    ("🤝","دعم محلي","فريق يفهمك ويتحدّث لغتك"),
]
xs=[Inches(6.95),Inches(0.85)]; y0=Inches(2.3)
for i,(e,t1,t2) in enumerate(points):
    col=i%2; row=i//2; x=xs[col]; y=y0+row*Inches(1.45)
    cc=rect(s,x,y,Inches(5.5),Inches(1.3),RGBColor(0x0A,0x3E,0x68),MSO_SHAPE.ROUNDED_RECTANGLE,lw=0); cc.adjustments[0]=0.08
    text(s,x+Inches(4.75),y+Inches(0.32),Inches(0.6),Inches(0.6),[[R(e,22,False,WHITE)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x+Inches(0.25),y+Inches(0.2),Inches(4.4),Inches(0.95),[[R(t1,16,True,WHITE)],[R(t2,12,False,RGBColor(0xBF,0xD8,0xEE))]],ls=1.05)
chip(s,Inches(5.0),Inches(6.75),Inches(3.3),Inches(0.6),AMBER)
text(s,Inches(5.0),Inches(6.75),Inches(3.3),Inches(0.6),[[R("جرّبه مجاناً اليوم",16,True,NAVY)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

prs.save("PayNest-Pricing.pptx")
print("SAVED PayNest-Pricing.pptx | slides:", len(prs.slides._sldIdLst))
